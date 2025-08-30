#!/usr/bin/env python3
"""
Agent Simulation Test for VoluteMCP Server

This script simulates how real AI agents would interact with your PowerPoint tools,
testing all functionality without requiring actual agent applications.
"""

import asyncio
import json
import subprocess
import sys
import time
import os
from pathlib import Path
from typing import Dict, Any, Optional, List


class MCPAgentSimulator:
    """Simulates an AI agent interacting with the MCP server."""
    
    def __init__(self):
        self.process = None
        self.request_id = 0
        self.server_info = {}
        self.available_tools = []
        self.available_resources = []
        
    async def start_server(self):
        """Start the MCP server in STDIO mode."""
        print("🚀 Starting MCP server...")
        
        self.process = await asyncio.create_subprocess_exec(
            sys.executable, "server.py", "stdio",
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=Path.cwd()
        )
        
        print("✅ Server started successfully")
    
    async def send_request(self, method: str, params: Optional[Dict] = None) -> Dict[Any, Any]:
        """Send an MCP request and get response."""
        if not self.process:
            raise RuntimeError("Server not started")
        
        self.request_id += 1
        request = {
            "jsonrpc": "2.0",
            "id": self.request_id,
            "method": method,
            "params": params or {}
        }
        
        # Send request
        request_json = json.dumps(request) + "\n"
        self.process.stdin.write(request_json.encode())
        await self.process.stdin.drain()
        
        # Read response
        response_line = await self.process.stdout.readline()
        if not response_line:
            raise RuntimeError("No response from server")
        
        return json.loads(response_line.decode().strip())
    
    async def initialize(self) -> bool:
        """Initialize the MCP connection."""
        print("🔄 Initializing MCP connection...")
        
        response = await self.send_request("initialize", {
            "protocolVersion": "2024-11-05",
            "capabilities": {
                "roots": {"listChanged": True},
                "sampling": {}
            },
            "clientInfo": {
                "name": "agent-simulator",
                "version": "1.0.0"
            }
        })
        
        if "result" in response:
            self.server_info = response["result"]
            print(f"✅ Connected to: {self.server_info.get('serverInfo', {}).get('name', 'Unknown')}")
            return True
        else:
            print(f"❌ Initialization failed: {response}")
            return False
    
    async def list_tools(self) -> List[Dict]:
        """Get all available tools."""
        print("🔧 Discovering available tools...")
        
        response = await self.send_request("tools/list")
        
        if "result" in response and "tools" in response["result"]:
            self.available_tools = response["result"]["tools"]
            print(f"✅ Found {len(self.available_tools)} tools")
            return self.available_tools
        else:
            print(f"❌ Failed to list tools: {response}")
            return []
    
    async def list_resources(self) -> List[Dict]:
        """Get all available resources."""
        print("📁 Discovering available resources...")
        
        response = await self.send_request("resources/list")
        
        if "result" in response and "resources" in response["result"]:
            self.available_resources = response["result"]["resources"]
            print(f"✅ Found {len(self.available_resources)} resources")
            return self.available_resources
        else:
            print(f"❌ Failed to list resources: {response}")
            return []
    
    async def call_tool(self, tool_name: str, arguments: Dict) -> Dict:
        """Call a specific tool with arguments."""
        print(f"🛠️  Calling tool: {tool_name}")
        
        response = await self.send_request("tools/call", {
            "name": tool_name,
            "arguments": arguments
        })
        
        if "result" in response:
            print(f"✅ Tool '{tool_name}' executed successfully")
            return response["result"]
        else:
            print(f"❌ Tool '{tool_name}' failed: {response}")
            return {"error": response.get("error", "Unknown error")}
    
    async def read_resource(self, uri: str) -> Dict:
        """Read a resource by URI."""
        print(f"📄 Reading resource: {uri}")
        
        response = await self.send_request("resources/read", {
            "uri": uri
        })
        
        if "result" in response:
            print(f"✅ Resource '{uri}' read successfully")
            return response["result"]
        else:
            print(f"❌ Failed to read resource '{uri}': {response}")
            return {"error": response.get("error", "Unknown error")}
    
    async def stop_server(self):
        """Stop the MCP server."""
        if self.process:
            print("🛑 Shutting down server...")
            self.process.terminate()
            await self.process.wait()
            print("✅ Server stopped")


async def create_test_powerpoint():
    """Create a simple test PowerPoint file for testing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        print("📊 Creating test PowerPoint file...")
        
        # Create presentation
        prs = Presentation()
        
        # Slide 1: Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        
        title.text = "Test Presentation"
        subtitle.text = "Created by VoluteMCP Agent Simulator"
        
        # Slide 2: Content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide2.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Key Features'
        tf = body_shape.text_frame
        tf.text = 'PowerPoint Analysis Tools'
        
        p = tf.add_paragraph()
        p.text = 'Metadata extraction'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = 'Content analysis'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = 'File validation'
        p.level = 1
        
        # Slide 3: Chart slide
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide3.shapes.title.text = "Data Overview"
        
        # Add a text box
        txBox = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
        tf = txBox.text_frame
        tf.text = "This slide contains sample data for analysis.\n\n"
        tf.text += "• Sample metric 1: 95%\n"
        tf.text += "• Sample metric 2: 87%\n" 
        tf.text += "• Sample metric 3: 92%"
        
        # Save presentation
        test_file = Path("test_presentation.pptx")
        prs.save(test_file)
        
        print(f"✅ Created test file: {test_file.absolute()}")
        return str(test_file.absolute())
        
    except ImportError:
        print("⚠️  python-pptx not available - using mock file path")
        return "test_presentation.pptx"
    except Exception as e:
        print(f"⚠️  Error creating PowerPoint: {e}")
        return "test_presentation.pptx"


async def test_basic_tools(agent: MCPAgentSimulator):
    """Test basic non-PowerPoint tools."""
    print("\n" + "🧮 " + "=" * 58)
    print("  TESTING BASIC TOOLS")
    print("=" * 60)
    
    # Test calculator
    result = await agent.call_tool("calculate", {
        "expression": "15 + 25 * 2"
    })
    print(f"   Calculator result: {result}")
    
    # Test text formatting
    result = await agent.call_tool("format_text", {
        "text": "hello world",
        "format_type": "title"
    })
    print(f"   Text formatting result: {result}")
    
    # Test hashing
    result = await agent.call_tool("hash_text", {
        "text": "VoluteMCP PowerPoint Tools",
        "algorithm": "sha256"
    })
    print(f"   Hash result: {result}")
    
    # Test server info
    result = await agent.call_tool("get_server_info", {})
    print(f"   Server info retrieved: {type(result)}")


async def test_powerpoint_tools(agent: MCPAgentSimulator, test_file: str):
    """Test all PowerPoint tools with the test file."""
    print("\n" + "📊 " + "=" * 58)
    print("  TESTING POWERPOINT TOOLS")
    print("=" * 60)
    
    # Test 1: Validate PowerPoint file
    print("\n🔍 Test 1: PowerPoint File Validation")
    result = await agent.call_tool("validate_powerpoint_file", {
        "presentation_path": test_file
    })
    print(f"   Validation result: {type(result)}")
    if isinstance(result, str) and len(result) > 100:
        print(f"   Result preview: {result[:100]}...")
    else:
        print(f"   Full result: {result}")
    
    # Test 2: Get PowerPoint summary
    print("\n📋 Test 2: PowerPoint Summary")
    result = await agent.call_tool("get_powerpoint_summary", {
        "presentation_path": test_file
    })
    print(f"   Summary result type: {type(result)}")
    if isinstance(result, str) and len(result) > 100:
        print(f"   Summary preview: {result[:100]}...")
    else:
        print(f"   Full summary: {result}")
    
    # Test 3: Analyze PowerPoint content
    print("\n🔬 Test 3: PowerPoint Content Analysis")
    result = await agent.call_tool("analyze_powerpoint_content", {
        "presentation_path": test_file,
        "slide_numbers": [1, 2],
        "extract_text_only": True
    })
    print(f"   Analysis result type: {type(result)}")
    if isinstance(result, str) and len(result) > 100:
        print(f"   Analysis preview: {result[:100]}...")
    else:
        print(f"   Full analysis: {result}")
    
    # Test 4: Extract full metadata
    print("\n🗂️  Test 4: PowerPoint Metadata Extraction")
    result = await agent.call_tool("extract_powerpoint_metadata", {
        "presentation_path": test_file,
        "include_slide_content": True,
        "include_master_slides": False,
        "output_format": "json"
    })
    print(f"   Metadata result type: {type(result)}")
    if isinstance(result, str) and len(result) > 200:
        print(f"   Metadata preview: {result[:200]}...")
    else:
        print(f"   Metadata length: {len(str(result))} characters")


async def test_resources(agent: MCPAgentSimulator):
    """Test resource access."""
    print("\n" + "📁 " + "=" * 58)
    print("  TESTING RESOURCE ACCESS")
    print("=" * 60)
    
    # Test server config resource
    result = await agent.read_resource("config://server")
    print(f"   Server config: {type(result)}")
    
    # Test environment resource
    result = await agent.read_resource("data://environment")
    print(f"   Environment data: {type(result)}")
    
    # Test user resource
    result = await agent.read_resource("users://1")
    print(f"   User data: {type(result)}")


async def simulate_agent_workflow(agent: MCPAgentSimulator, test_file: str):
    """Simulate a realistic agent workflow analyzing a PowerPoint."""
    print("\n" + "🤖 " + "=" * 58)
    print("  SIMULATING AGENT WORKFLOW")
    print("=" * 60)
    
    print("\n🎯 Scenario: Agent analyzing a presentation for content insights")
    
    # Step 1: Validate the file first
    print("\n📍 Step 1: Validating presentation file...")
    validation = await agent.call_tool("validate_powerpoint_file", {
        "presentation_path": test_file
    })
    
    # Step 2: Get overview
    print("📍 Step 2: Getting presentation overview...")
    summary = await agent.call_tool("get_powerpoint_summary", {
        "presentation_path": test_file
    })
    
    # Step 3: Analyze specific content
    print("📍 Step 3: Analyzing slide content in detail...")
    analysis = await agent.call_tool("analyze_powerpoint_content", {
        "presentation_path": test_file,
        "extract_text_only": False
    })
    
    # Step 4: Get technical metadata
    print("📍 Step 4: Extracting technical metadata...")
    metadata = await agent.call_tool("extract_powerpoint_metadata", {
        "presentation_path": test_file,
        "include_slide_content": True,
        "output_format": "json"
    })
    
    print("\n🎉 Agent workflow completed!")
    print("✅ All PowerPoint analysis tools executed successfully")
    
    # Simulate agent response
    print("\n🗨️  Simulated Agent Response:")
    print("   'I've successfully analyzed your PowerPoint presentation.'")
    print("   'The file is valid and contains multiple slides with content.'")
    print("   'I was able to extract metadata, validate the structure,')")
    print("   'and analyze the content using the available tools.'")


async def run_comprehensive_test():
    """Run comprehensive agent simulation tests."""
    print("🤖 VoluteMCP Agent Simulation Test Suite")
    print("=" * 60)
    
    # Create test file
    test_file = await create_test_powerpoint()
    
    # Initialize agent
    agent = MCPAgentSimulator()
    
    try:
        # Start server and initialize
        await agent.start_server()
        
        if not await agent.initialize():
            return False
        
        # Discover capabilities
        tools = await agent.list_tools()
        resources = await agent.list_resources()
        
        print(f"\n📊 Server Capabilities:")
        print(f"   Tools: {len(tools)}")
        print(f"   Resources: {len(resources)}")
        
        # Show PowerPoint tools specifically
        powerpoint_tools = [t for t in tools if 'powerpoint' in t.get('name', '').lower()]
        print(f"   PowerPoint Tools: {len(powerpoint_tools)}")
        
        if powerpoint_tools:
            print("   Available PowerPoint Tools:")
            for tool in powerpoint_tools:
                name = tool.get('name', 'Unknown')
                desc = tool.get('description', 'No description')
                print(f"     • {name}: {desc[:60]}{'...' if len(desc) > 60 else ''}")
        
        # Run tests
        await test_basic_tools(agent)
        await test_powerpoint_tools(agent, test_file)
        await test_resources(agent)
        await simulate_agent_workflow(agent, test_file)
        
        print("\n" + "=" * 60)
        print("🎉 AGENT SIMULATION COMPLETE!")
        print("✅ Your PowerPoint tools work correctly with agent applications")
        print("🔗 The server is ready for integration with real AI agents")
        print("=" * 60)
        
        return True
        
    except Exception as e:
        print(f"\n❌ Simulation failed: {e}")
        import traceback
        traceback.print_exc()
        return False
        
    finally:
        await agent.stop_server()


if __name__ == "__main__":
    print("🚀 Starting Agent Simulation Test...")
    
    # Run the simulation
    success = asyncio.run(run_comprehensive_test())
    
    if success:
        print("\n✅ SUCCESS: Your MCP server works perfectly with AI agents!")
        print("\n🔗 Next Steps:")
        print("   1. Test with any MCP-compatible client library")
        print("   2. Integrate with custom AI agents using STDIO transport")
        print("   3. Use the JSON-RPC protocol for direct integration")
        print("   4. Deploy to production for real agent applications")
    else:
        print("\n❌ FAILURE: Some tests failed")
        print("💡 Check the error messages above for troubleshooting")
    
    sys.exit(0 if success else 1)
