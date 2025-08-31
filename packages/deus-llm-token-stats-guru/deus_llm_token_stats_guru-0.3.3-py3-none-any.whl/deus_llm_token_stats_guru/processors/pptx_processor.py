"""PPTX (PowerPoint) file processor."""

from pathlib import Path
from typing import Set, List
import logging

from .base import BaseFileProcessor
from ..types import CountResult
from .._exceptions import FileProcessingError

logger = logging.getLogger(__name__)


class PPTXProcessor(BaseFileProcessor):
    """Processor for PowerPoint PPTX files."""
    
    @property
    def supported_extensions(self) -> Set[str]:
        """Return supported file extensions."""
        return {".pptx", ".ppt"}
    
    @property
    def processor_name(self) -> str:
        """Return processor name."""
        return "PowerPoint"
    
    def extract_text(self, file_path: Path) -> List[str]:
        """Extract text content from PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            List of text strings from all slides
            
        Raises:
            FileProcessingError: If PPTX cannot be read
        """
        try:
            # Try python-pptx first (most comprehensive for .pptx)
            try:
                from pptx import Presentation
                return self._extract_with_python_pptx(file_path)
            except ImportError:
                logger.debug("python-pptx not available")
                pass
            
            # Fallback to extraction via zip (PPTX is a zip file)
            try:
                import zipfile
                import xml.etree.ElementTree as ET
                return self._extract_with_zip(file_path)
            except Exception as e:
                logger.debug(f"Zip extraction failed: {e}")
                pass
            
            raise FileProcessingError(
                f"No PPTX processing libraries available. Install: pip install python-pptx"
            )
            
        except Exception as e:
            logger.warning(f"⚠️  Failed to extract text from {file_path}: {e}")
            return []
    
    def _extract_with_python_pptx(self, file_path: Path) -> List[str]:
        """Extract text using python-pptx library."""
        from pptx import Presentation
        
        text_content = []
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"Slide {slide_num}"]
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and len(text) > 1:
                        slide_text.append(text)
                
                # Handle tables in slides
                if hasattr(shape, "table"):
                    table_text = self._extract_table_text(shape.table)
                    if table_text:
                        slide_text.extend(table_text)
                
                # Handle grouped shapes
                if hasattr(shape, "shapes"):
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text") and sub_shape.text:
                            text = sub_shape.text.strip()
                            if text and len(text) > 1:
                                slide_text.append(text)
            
            # Extract notes if available
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text and len(notes_text) > 1:
                    slide_text.append(f"Notes: {notes_text}")
            
            if len(slide_text) > 1:  # More than just slide number
                text_content.extend(slide_text)
        
        return text_content
    
    def _extract_table_text(self, table) -> List[str]:
        """Extract text from PowerPoint table."""
        table_text = []
        
        try:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text:
                        cell_content = cell.text.strip()
                        if cell_content and len(cell_content) > 1:
                            row_text.append(cell_content)
                
                if row_text:
                    table_text.append(" ".join(row_text))
        
        except Exception as e:
            logger.debug(f"⚠️  Failed to extract table text: {e}")
        
        return table_text
    
    def _extract_with_zip(self, file_path: Path) -> List[str]:
        """Extract text using direct ZIP parsing (fallback method)."""
        import zipfile
        import xml.etree.ElementTree as ET
        
        text_content = []
        
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Get list of slide files
                slide_files = [f for f in zip_file.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                slide_files.sort()  # Ensure proper order
                
                for slide_num, slide_file in enumerate(slide_files, 1):
                    slide_text = [f"Slide {slide_num}"]
                    
                    try:
                        with zip_file.open(slide_file) as xml_file:
                            tree = ET.parse(xml_file)
                            root = tree.getroot()
                            
                            # Extract text from all text elements
                            # PowerPoint XML uses various namespaces
                            for elem in root.iter():
                                if elem.text and elem.text.strip():
                                    text = elem.text.strip()
                                    if len(text) > 1 and text not in slide_text:
                                        slide_text.append(text)
                    
                    except Exception as e:
                        logger.debug(f"⚠️  Failed to parse slide {slide_file}: {e}")
                        continue
                    
                    if len(slide_text) > 1:
                        text_content.extend(slide_text)
        
        except Exception as e:
            logger.debug(f"⚠️  ZIP extraction failed: {e}")
        
        return text_content
    
    def process_file(self, file_path: Path) -> CountResult:
        """Process PPTX file with enhanced metadata.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            CountResult with PowerPoint-specific statistics
        """
        logger.info(f"🎞️  Processing PowerPoint file: {file_path.name}")
        
        # Extract text content
        text_content = self.extract_text(file_path)
        total_tokens = sum(self.count_tokens_in_text(text) for text in text_content)
        
        # Get basic file statistics
        file_size = file_path.stat().st_size
        
        # Estimate slides and content items from text
        slide_count = len([line for line in text_content if line.startswith("Slide ")])
        content_items = len([line for line in text_content if not line.startswith("Slide ")])
        
        result: CountResult = {
            "file_path": str(file_path),
            "total_tokens": total_tokens,
            "row_count": content_items,  # Content items (text boxes, etc.)
            "column_count": slide_count,  # Number of slides
            "encoding_model": self.encoding_model,
            "file_size_bytes": file_size,
        }
        
        logger.info(f"✅ Processed {file_path.name}: {total_tokens:,} tokens, {content_items} items, {slide_count} slides")
        return result