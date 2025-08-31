#!/usr/bin/env python3

from fastmcp import FastMCP
from pydantic import Field, BaseModel
from pptx import Presentation
import os

mcp = FastMCP("Generate PPTs for MM PTMs easily")

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = {
    "math": os.path.join(SCRIPT_DIR, "template.pptx"),
    "science": os.path.join(SCRIPT_DIR, "Science_template.pptx")
}
OUTPUT_FOLDER = os.path.join(SCRIPT_DIR, "generated_ppts")
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


class MathStudent(BaseModel):
    #  general details
    Tutor: str = Field(..., description="Tutor's name:")
    Student: str = Field(..., description="Student's name:")
    Subjects: str = Field(..., description="Student's subject:")
    ParentRequirement: str = Field(..., description="Parent Requirement:")
    ReportingPeriod: str = Field(..., description="Reoprting period:")
    NoOfSessions: int = Field(..., description="No of sessions:")
    # IXL Diagnostic Stats
    Target: int = Field(..., description="IXL Target score")
    Numbers: int = Field(..., description="IXL Numbers and Operations Score")
    Algebra: int = Field(..., description="IXL Algebra and Algebraic Thinking Score")
    Fractions: int = Field(..., description="IXL Fractions Score")
    Geometry: int = Field(..., description="IXL Goemetry Score")
    Measurement: int = Field(..., description="IXL Measurement Score")
    Data: int = Field(..., description="IXL Data & Probability Score")
    Overall: int = Field(..., description="Overall IXL Math Level")
    #  add the fields for the third slide here to talk about the rec skills from IXL
    IXLAreaOfImprovement1: str = Field(..., description="Suggested IXL strand to work on:")
    IXLAreaOfImprovement2: str = Field(..., description="IXL Suggested strand 2 to work on:")
    AreaOfImprovement1SuggestedSkill1: str = Field(..., description="Skill 1 to work on to improve area of improvement 1")
    AreaOfImprovement1SuggestedSkill2: str = Field(..., description="Skill 2 to work on to improve area of improvement 1")
    AreaOfImprovement2SuggestedSkill1: str = Field(..., description="Skill 1 to work on to improve area of improvement 2")
    AreaOfImprovement2SuggestedSkill2: str = Field(..., description="Skill 2 to work on to improve area of improvement 2")
    # topics taught that month
    Topic1: str = Field(..., description="Topic 1 covered this month")
    T1Status: str = Field(..., description="Status of the topic 1")
    Topic2: str = Field(..., description="Topic 2 covered this month")
    T2Status: str = Field(..., description="Status of the topic 2")
    # monthly test details
    MTest: int = Field(..., description="Monthly Test Score (out of 25)")
    # learning gaps and action plan
    LGap: str = Field(..., description="Learning gap identified")
    APlan: str = Field(..., description="Action Plan for the learning Gap")
    StudentStepsNeeded: str = Field(..., description="Steps needed from student")
    # upcoming tasks
    Task1: str = Field(..., description="Next task planned")
    Task1Sess: str = Field(..., description="Number of sessions needed for task 1")
    Task2: str = Field(..., description="Next task2 planned")
    Task2Sess: str = Field(..., description="Number of sessions needed for task 2")
    Notes: str = Field(..., description="Notes for upcoming tasks planned")

class ScienceStudent(BaseModel):
    Tutor: str = Field(..., description="Tutor's name:")
    Student: str = Field(..., description="Student's name:")
    Subjects: str = Field(..., description="Student's subject:")
    ParentRequirement: str = Field(..., description="Parent Requirement:")
    ReportingPeriod: str = Field(..., description="Reoprting period:")
    NoOfSessions: int = Field(..., description="No of sessions:")
    # topics taught that month
    Topic1: str = Field(..., description="Topic 1 covered this month")
    T1Status: str = Field(..., description="Status of the topic 1")
    Topic2: str = Field(..., description="Topic 2 covered this month")
    T2Status: str = Field(..., description="Status of the topic 2")
    # monthly test details
    MTest: int = Field(..., description="Monthly Test Score (out of 25)")
    # learning gaps and action plan
    LGap: str = Field(..., description="Learning gap identified")
    APlan: str = Field(..., description="Action Plan for the learning Gap")
    StudentStepsNeeded: str = Field(..., description="Steps needed from student")
    # upcoming tasks
    Task1: str = Field(..., description="Next task planned")
    Task1Sess: str = Field(..., description="Number of sessions needed for task 1")
    Task2: str = Field(..., description="Next task2 planned")
    Task2Sess: str = Field(..., description="Number of sessions needed for task 2")
    Notes: str = Field(..., description="Notes for upcoming tasks planned")





def replace_text(shape, data: dict):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in data.items():
                    placeholder = f"{{{{{key}}}}}"
                    if placeholder in run.text:
                        if key == "Subjects":
                            value = ", ".join([s.strip() for s in value.split(",")])
                        run.text = run.text.replace(placeholder, str(value))
    if shape.shape_type == 6:
        for subshape in shape.shapes:
            replace_text(subshape, data)
    


    if shape.shape_type == 19:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in data.items():
                            placeholder= f"{{{{{key}}}}}"
                            if placeholder in run.text:
                                if key == "Subjects":
                                    value = ", ".join([s.strip() for s in value.split(",")])
                                run.text = run.text.replace(placeholder, str(value))

def add_stuff_to_ppt(prs: Presentation, data: BaseModel):
    d = data.dict()
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text(shape, d)



@mcp.tool
def math_ppt(student_type:str, student_data:dict):
    """Generate PPTs for MM PTMs easily"""
    student = MathStudent(**student_data)
    prs = Presentation(TEMPLATE_PATH["math"])
    add_stuff_to_ppt(prs, student)
    output_path = os.path.join(OUTPUT_FOLDER, f"{student.Student}_{student_type}.pptx")
    prs.save(output_path)
    return f"PPT generated for {student.Student} at {output_path}"

@mcp.tool
def science_ppt(student_type:str, student_data:dict):
    """Generate PPTs for MM PTMs easily"""
    student = ScienceStudent(**student_data)
    prs = Presentation(TEMPLATE_PATH["science"])
    add_stuff_to_ppt(prs, student)
    output_path = os.path.join(OUTPUT_FOLDER, f"{student.Student}_{student_type}.pptx")
    prs.save(output_path)
    return f"PPT generated for {student.Student} at {output_path}"



def main():
    mcp.run()

if __name__=="__main__":
    main()

