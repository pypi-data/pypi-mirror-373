# split_txts.py
import argparse
from io import BytesIO
import tempfile
from typing import Optional, List

import chardet
import fitz  # PyMuPDF
import os
import pandas as pd
import re
import subprocess
from subprocess import TimeoutExpired
from docx import Document
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from botrun_ask_folder.drive_download import (
    append_export_extension_to_path,
    truncate_filename,
)
from botrun_ask_folder.models.drive_file import DriveFile
from botrun_ask_folder.models.rag_metadata import RagMetadata, PAGE_NUMBER_NA
from botrun_ask_folder.services.storage.storage_client import StorageClient
from botrun_ask_folder.services.storage.storage_factory import storage_client_factory
from .drive_download_metadata import (
    update_download_metadata_after_process_file,
    get_drive_download_metadata,
)
from .emoji_progress_bar import EmojiProgressBar
from .is_cached import is_cached
from .is_text_file import is_text_file
from .log_handler import LogHandler
from .safe_join import safe_join
from pptx.exc import PackageNotFoundError

log_handler = LogHandler("./logs", "split_txts.py")  # 創建日誌處理器實例


def get_page_number_prefix():
    return "page_"  # f"page_{page_num}.txt"


def get_part_number_prefix():
    return "part_"  # f"page_{page_num}.txt"


def get_file_name_with_page_number(file_name: str, page_num: int) -> str:
    return f"{file_name}.{get_page_number_prefix()}{page_num}.txt"


def re_compile_page_number_pattern():
    return re.compile(rf"\.{get_page_number_prefix()}\d+\.txt$")


def get_page_number(str_file_path):
    pattern = re.compile(rf"\.{get_page_number_prefix()}(\d+)\.txt$")
    match = pattern.search(str_file_path)
    if match:
        return int(match.group(1))
    else:
        return 0


def write_to_file(
    output_folder: str,
    filename: str,
    text: str,
    chars_per_page: int,
    add_page_numbers: bool = True,
):
    output_path = Path(output_folder)
    # pdf like files, 如果 add_page_numbers 為 False，則不加入頁碼，並且完整寫入檔案不要切割，此時 chars_per_page 無用
    if not add_page_numbers:
        output_filename = f"{filename}.txt"
        with open(output_path / output_filename, "w", encoding="utf-8") as file:
            file.write(text)
        return
    # txt like files, 需要切割, 按照 chars_per_page 大小切割
    for page_num, i in enumerate(range(0, len(text), chars_per_page), start=1):
        page_filename = get_file_name_with_page_number(filename, page_num)
        with open(output_path / page_filename, "w", encoding="utf-8") as file:
            file.write(text[i : i + chars_per_page])


def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    return " ".join(text)


def extract_text_from_pptx(slide):
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(shape.text)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        text.append(cell.text_frame.text)
    return " ".join(text)


def process_powerpoint(
    max_chars_per_file: int,
    file_name: str,
    file_type: str,
    output_folder: str,
    source_file: str,
    gen_page_imgs: bool = False,
) -> List[RagMetadata]:
    # 确保文件是 .pptx 格式
    expected_pptx_path = safe_join(output_folder, f"{file_name}.pptx")
    if not os.path.exists(expected_pptx_path) or file_type == ".ppt":
        if file_type == ".ppt":
            try:
                converted_file_path = convert_office_file(source_file, file_type)
            except Exception as e:
                print(f"Failed to convert .ppt file: {e}")
                return []
        else:
            converted_file_path = source_file
    else:
        converted_file_path = expected_pptx_path

    try:
        prs = Presentation(converted_file_path)
    except PackageNotFoundError:
        print(f"PowerPoint file not found: {converted_file_path}")
        return []
    except Exception as e:
        print(f"Error opening PowerPoint file: {e}")
        return []

    lst_rag_metadata = []

    for slide_number, slide in enumerate(prs.slides):
        try:
            text = extract_text_from_pptx(slide)
        except Exception as e:
            print(f"Error extracting text from slide {slide_number + 1}: {e}")
            continue

        # 将幻灯片文本分割成段落
        paragraphs = text.split("\n")

        current_chunk = []
        current_char_count = 0
        chunk_number = 1

        for paragraph in paragraphs:
            if (
                current_char_count + len(paragraph) > max_chars_per_file
                and current_chunk
            ):
                # 储存当前块并开始新的块
                try:
                    save_chunk(
                        current_chunk,
                        output_folder,
                        file_name,
                        f"{get_page_number_prefix()}{slide_number + 1}",
                        f"{chunk_number:03d}",
                    )
                    lst_rag_metadata.append(
                        RagMetadata(
                            name=get_save_chunk_file_name(
                                file_name,
                                f"{get_page_number_prefix()}{slide_number + 1}",
                                f"{chunk_number:03d}",
                            ),
                            ori_file_name=f"{file_name}{file_type}",
                            gen_page_imgs=gen_page_imgs,
                            page_number=slide_number + 1,
                            save_path=get_save_chunk_file_path(
                                output_folder,
                                file_name,
                                f"{get_page_number_prefix()}{slide_number + 1}",
                                f"{chunk_number:03d}",
                            ),
                        )
                    )
                except Exception as e:
                    print(
                        f"Error saving chunk for slide {slide_number + 1}, chunk {chunk_number}: {e}"
                    )

                current_chunk = []
                current_char_count = 0
                chunk_number += 1

            current_chunk.append(paragraph)
            current_char_count += len(paragraph)

        # 储存最后一个块（如果有的话）
        if current_chunk:
            try:
                save_chunk(
                    current_chunk,
                    output_folder,
                    file_name,
                    f"{get_page_number_prefix()}{slide_number + 1}",
                    f"{chunk_number:03d}",
                )
                lst_rag_metadata.append(
                    RagMetadata(
                        name=get_save_chunk_file_name(
                            file_name,
                            f"{get_page_number_prefix()}{slide_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                        ori_file_name=f"{file_name}{file_type}",
                        gen_page_imgs=gen_page_imgs,
                        page_number=slide_number + 1,
                        save_path=get_save_chunk_file_path(
                            output_folder,
                            file_name,
                            f"{get_page_number_prefix()}{slide_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                    )
                )
            except Exception as e:
                print(f"Error saving last chunk for slide {slide_number + 1}: {e}")

    return lst_rag_metadata


def process_excel(
    original_file_path: str,
    output_folder: str,
    file_name: str,
    file_type: str,
    max_chars_per_file: int = 5000,
) -> List[RagMetadata]:
    # seba max token 是 8192，但是 max_chars 不見得會剛好 meet 到該 token 數，所以中間抓一個 buffer
    if file_type == ".xlsx":
        xls = pd.ExcelFile(original_file_path, engine="openpyxl")
    elif file_type == ".xls":
        xls = pd.ExcelFile(original_file_path, engine="xlrd")
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    lst_rag_metadata = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name)

        # 將 DataFrame 轉換為字符串列表
        content_list = df.apply(lambda row: " ".join(row.astype(str)), axis=1).tolist()

        current_chunk = []
        current_char_count = 0
        file_counter = 1

        for row in content_list:
            if current_char_count + len(row) > max_chars_per_file and current_chunk:
                # 儲存當前塊並開始新的塊
                save_chunk(
                    current_chunk,
                    output_folder,
                    file_name,
                    f"{sheet_name}_{get_page_number_prefix()}_1_",
                    file_counter,
                )
                lst_rag_metadata.append(
                    RagMetadata(
                        name=get_save_chunk_file_name(
                            file_name,
                            f"{sheet_name}_{get_page_number_prefix()}_1_",
                            file_counter,
                        ),
                        ori_file_name=f"{file_name}{file_type}",
                        gen_page_imgs=False,
                        page_number=PAGE_NUMBER_NA,
                        sheet_name=sheet_name,
                        save_path=get_save_chunk_file_path(
                            output_folder,
                            file_name,
                            f"{sheet_name}_{get_page_number_prefix()}_1_",
                            file_counter,
                        ),
                    )
                )
                current_chunk = []
                current_char_count = 0
                file_counter += 1

            current_chunk.append(row)
            current_char_count += len(row)

        # 儲存最後一個塊（如果有的話）
        if current_chunk:
            save_chunk(
                current_chunk,
                output_folder,
                file_name,
                f"{sheet_name}_{get_page_number_prefix()}_1_",
                file_counter,
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        file_name,
                        f"{sheet_name}_{get_page_number_prefix()}_1_",
                        file_counter,
                    ),
                    ori_file_name=f"{file_name}{file_type}",
                    gen_page_imgs=False,
                    page_number=PAGE_NUMBER_NA,
                    sheet_name=sheet_name,
                    save_path=get_save_chunk_file_path(
                        output_folder,
                        file_name,
                        f"{sheet_name}_{get_page_number_prefix()}_1_",
                        file_counter,
                    ),
                )
            )
    return lst_rag_metadata


def get_save_chunk_file_name(file_name, sheet_name, file_counter):
    return f"{file_name}_{sheet_name}.{get_part_number_prefix()}{file_counter}.txt"


def get_save_chunk_file_path(output_folder, file_name, sheet_name, file_counter):
    output_file_name = get_save_chunk_file_name(file_name, sheet_name, file_counter)
    output_path = os.path.join(output_folder, output_file_name)
    return output_path


def get_save_chunk_file_path_by_storage(
    output_folder, file_name, sheet_name, file_counter
):
    output_file_name = get_save_chunk_file_name(file_name, sheet_name, file_counter)
    output_path = "/".join([output_folder, output_file_name])
    return output_path


def save_chunk(chunk, output_folder, file_name, sheet_name, file_counter):
    # output_file_name = get_save_chunk_file_name(file_name, sheet_name, file_counter)
    output_path = get_save_chunk_file_path(
        output_folder, file_name, sheet_name, file_counter
    )

    with open(output_path, "w", encoding="utf-8") as f:
        for row in chunk:
            f.write(row + "\n")


def convert_office_file(file_path: str, target_extension: str):
    converted_file_dir = os.path.dirname(file_path)
    converted_file_base = os.path.splitext(os.path.basename(file_path))[0]
    converted_file_name = f"{converted_file_base}.{target_extension.lstrip('.')}"
    converted_file_path = safe_join(converted_file_dir, converted_file_name)

    # 建構 LibreOffice 的命令列格式，用於檔案轉換
    libreoffice_command = os.getenv("LIBRE_OFFICE_COMMAND", "libreoffice")
    cmd = [
        libreoffice_command,
        "--headless",
        "--convert-to",
        target_extension.lstrip("."),
        "--outdir",
        converted_file_dir,
        file_path,
    ]
    try:
        # 設置 5 分鐘 (300 秒) 的超時限制
        proc = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            check=True,
            timeout=180,  # 3 分鐘超時
        )
        stdout, stderr = proc.stdout, proc.stderr
        log_handler.write_log(f"Stdout: {stdout}")
        log_handler.write_log(f"Stderr: {stderr}")
    except TimeoutExpired:
        error_message = f"Conversion of {file_path} timed out after 3 minutes"
        log_handler.write_log(error_message)
        raise Exception(error_message)
    except subprocess.CalledProcessError as e:
        error_message = f"Error converting {file_path}. Return code: {e.returncode}\nStdout: {e.stdout}\nStderr: {e.stderr}"
        log_handler.write_log(error_message)
        raise Exception(error_message)
    except Exception as e:
        log_handler.write_log(f"Error converting {file_path}: {e}")
        raise Exception(f"Error converting {file_path}: {e}")

    if not os.path.exists(converted_file_path):
        error_message = f"Converted file not found: {converted_file_path}"
        log_handler.write_log(error_message)
        raise Exception(error_message)

    return converted_file_path


def process_single_file(
    source_file: str,
    output_folder: str,
    chars_per_page: int,
    force: bool,
    stats: dict = None,
    gen_page_imgs: bool = False,
) -> List[RagMetadata]:
    """
    @param metadata_dir: str, metadata 資料夾路徑，有傳才會更新 metadata
    """
    file_name, file_type = os.path.splitext(os.path.basename(source_file))
    if is_cached(output_folder, force) and stats is not None:
        stats["cached"] += 1
        return []
    os.makedirs(output_folder, exist_ok=True)
    initial_files_count = len(list(Path(output_folder).glob("*")))
    lst_rag_metadata = []
    if file_type in [".docx", ".doc"]:
        lst_rag_metadata = process_word(
            chars_per_page, file_name, file_type, output_folder, source_file
        )
    if file_type == ".pdf":
        lst_rag_metadata = process_pdf(
            chars_per_page,
            file_name,
            file_type,
            output_folder,
            source_file,
            gen_page_imgs,
        )
    if file_type in [".ppt", ".pptx"]:
        lst_rag_metadata = process_powerpoint(
            chars_per_page,
            file_name,
            file_type,
            output_folder,
            source_file,
            gen_page_imgs=False,
        )
    if file_type in [".xls", ".xlsx"]:
        lst_rag_metadata = process_excel(
            source_file, output_folder, file_name, file_type
        )
    if file_type == ".ods":
        converted_file_path = convert_office_file(source_file, ".xlsx")
        lst_rag_metadata = process_excel(
            converted_file_path, output_folder, file_name, ".xlsx"
        )
        os.remove(converted_file_path)
    if file_type == ".odp":
        converted_file_path = convert_office_file(source_file, ".pptx")
        lst_rag_metadata = process_powerpoint(
            chars_per_page,
            file_name,
            file_type,
            output_folder,
            source_file,
            gen_page_imgs=False,
        )
        os.remove(converted_file_path)
    if file_type == ".rtf" or file_type == ".odt":
        converted_file_path = convert_office_file(source_file, ".docx")
        lst_rag_metadata = process_word(
            chars_per_page, file_name, ".docx", output_folder, converted_file_path
        )
        os.remove(converted_file_path)
    if file_type == ".txt" or is_text_file(source_file):
        lst_rag_metadata = process_txt(
            chars_per_page, file_name, file_type, output_folder, source_file
        )
    final_files_count = len(list(Path(output_folder).glob("*")))
    if stats is not None:
        stats["generated"] += final_files_count - initial_files_count
    return lst_rag_metadata


def process_txt(
    max_chars_per_file: int,
    file_name: str,
    file_type: str,
    output_folder: str,
    source_file: str,
) -> List[RagMetadata]:
    # 檢測文件編碼
    with open(source_file, "rb") as f:
        raw_data = f.read()
    detected = chardet.detect(raw_data)
    encoding = detected["encoding"]

    file_counter = 1
    current_chunk = []
    current_char_count = 0
    lst_rag_metadata = []

    def save_current_chunk():
        nonlocal file_counter, current_chunk, current_char_count, lst_rag_metadata
        if current_chunk:
            content = "\n".join(current_chunk)
            save_chunk(
                [content],
                output_folder,
                file_name,
                f"{get_page_number_prefix()}_1_",
                f"{file_counter:03d}",
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        file_name,
                        f"{get_page_number_prefix()}_1_",
                        f"{file_counter:03d}",
                    ),
                    ori_file_name=f"{file_name}{file_type}",
                    gen_page_imgs=False,
                    page_number=PAGE_NUMBER_NA,
                    sheet_name="",
                    save_path=get_save_chunk_file_path(
                        output_folder,
                        file_name,
                        f"{get_page_number_prefix()}_1_",
                        f"{file_counter:03d}",
                    ),
                )
            )
            current_chunk = []
            current_char_count = 0
            file_counter += 1

    def split_line(line, max_length):
        parts = []
        current_part = ""
        for char in line:
            if len(current_part) + len(char) > max_length:
                parts.append(current_part)
                current_part = char
            else:
                current_part += char
        if current_part:
            parts.append(current_part)
        return parts

    with open(source_file, "r", encoding=encoding, errors="replace") as f:
        for line in f:
            line = line.strip()
            if not line:
                if current_chunk and current_char_count + 1 <= max_chars_per_file:
                    current_chunk.append("")
                    current_char_count += 1
                continue

            line_parts = split_line(line, max_chars_per_file)
            for part in line_parts:
                if current_char_count + len(part) + 1 > max_chars_per_file:
                    save_current_chunk()

                current_chunk.append(part)
                current_char_count += len(part) + 1  # +1 for newline

                if current_char_count >= max_chars_per_file:
                    save_current_chunk()

    if current_chunk:
        save_current_chunk()  # 保存最後的內容

    return lst_rag_metadata


def process_pdf(
    max_chars_per_file: int,
    file_name: str,
    file_type: str,
    output_folder: str,
    source_file: str,
    gen_page_imgs: bool = False,
) -> List[RagMetadata]:
    doc = fitz.open(source_file)
    lst_rag_metadata = []

    for page_number in range(len(doc)):
        page = doc.load_page(page_number)
        text = page.get_text()

        # 將頁面文本分割成段落
        paragraphs = text.split("\n")

        current_chunk = []
        current_char_count = 0
        chunk_number = 1

        for paragraph in paragraphs:
            if (
                current_char_count + len(paragraph) > max_chars_per_file
                and current_chunk
            ):
                # 儲存當前塊並開始新的塊
                save_chunk(
                    current_chunk,
                    output_folder,
                    file_name,
                    f"{get_page_number_prefix()}{page_number + 1}",
                    f"{chunk_number:03d}",
                )
                lst_rag_metadata.append(
                    RagMetadata(
                        name=get_save_chunk_file_name(
                            file_name,
                            f"{get_page_number_prefix()}{page_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                        ori_file_name=f"{file_name}{file_type}",
                        gen_page_imgs=gen_page_imgs,
                        page_number=page_number + 1,
                        save_path=get_save_chunk_file_path(
                            output_folder,
                            file_name,
                            f"{get_page_number_prefix()}{page_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                    )
                )
                current_chunk = []
                current_char_count = 0
                chunk_number += 1

            current_chunk.append(paragraph)
            current_char_count += len(paragraph)

        # 儲存最後一個塊（如果有的話）
        if current_chunk:
            save_chunk(
                current_chunk,
                output_folder,
                file_name,
                f"{get_page_number_prefix()}{page_number + 1}",
                f"{chunk_number:03d}",
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        file_name,
                        f"{get_page_number_prefix()}{page_number + 1}",
                        f"{chunk_number:03d}",
                    ),
                    ori_file_name=f"{file_name}{file_type}",
                    gen_page_imgs=gen_page_imgs,
                    page_number=page_number + 1,
                    save_path=get_save_chunk_file_path(
                        output_folder,
                        file_name,
                        f"{get_page_number_prefix()}{page_number + 1}",
                        f"{chunk_number:03d}",
                    ),
                )
            )

    return lst_rag_metadata


def process_word(
    chars_per_page, file_name, file_type, output_folder, source_file
) -> List[RagMetadata]:
    try:
        pdf_path = convert_office_file(source_file, ".pdf")
        lst_rag_metadata = process_pdf(
            chars_per_page, file_name, ".pdf", output_folder, pdf_path
        )
        os.remove(pdf_path)
    except Exception as e:
        log_handler.write_log(
            f"Conversion of {source_file} to PDF {str(e)}. Processing Word file directly."
        )
        lst_rag_metadata = process_word_directly(
            chars_per_page, file_name, file_type, output_folder, source_file
        )
    return lst_rag_metadata


def process_word_directly(
    chars_per_page, file_name, file_type, output_folder, source_file
) -> List[RagMetadata]:
    doc = Document(source_file)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    text = "\n".join(text)

    paragraphs = text.split("\n")
    current_chunk = []
    current_char_count = 0
    chunk_number = 1
    lst_rag_metadata = []

    for paragraph in paragraphs:
        if current_char_count + len(paragraph) > chars_per_page and current_chunk:
            save_chunk(
                current_chunk,
                output_folder,
                file_name,
                f"{get_page_number_prefix()}1",
                f"{chunk_number:03d}",
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        file_name,
                        f"{get_page_number_prefix()}1",
                        f"{chunk_number:03d}",
                    ),
                    ori_file_name=f"{file_name}{file_type}",
                    gen_page_imgs=False,
                    page_number=PAGE_NUMBER_NA,
                    save_path=get_save_chunk_file_path(
                        output_folder,
                        file_name,
                        f"{get_page_number_prefix()}1",
                        f"{chunk_number:03d}",
                    ),
                )
            )
            current_chunk = []
            current_char_count = 0
            chunk_number += 1

        current_chunk.append(paragraph)
        current_char_count += len(paragraph)

    if current_chunk:
        save_chunk(
            current_chunk,
            output_folder,
            file_name,
            f"{get_page_number_prefix()}1",
            f"{chunk_number:03d}",
        )
        lst_rag_metadata.append(
            RagMetadata(
                name=get_save_chunk_file_name(
                    file_name,
                    f"{get_page_number_prefix()}1",
                    f"{chunk_number:03d}",
                ),
                ori_file_name=f"{file_name}{file_type}",
                gen_page_imgs=False,
                page_number=PAGE_NUMBER_NA,
                save_path=get_save_chunk_file_path(
                    output_folder,
                    file_name,
                    f"{get_page_number_prefix()}1",
                    f"{chunk_number:03d}",
                ),
            )
        )

    return lst_rag_metadata


def split_txts_no_threads(
    source_files: list,
    output_folders: list,
    chars_per_page: int = 2000,
    force: bool = False,
    gen_page_imgs: bool = False,
    metadata_dir: str = None,
    new_items: list = None,
    google_drive_folder_id: str = None,
):
    """
    @param metadata_dir: str, metadata 資料夾路徑，有傳才會更新 metadata
    @param new_items: 這裡的 item的格式，是drive list 出來的結果，可以參考 metadata，如果是 None，執行的方式會照舊(目前應該是會有 bug，但是為不影響其它用到的地方)，如果有傳，包含空 list，會檢查在這個 list 裡的檔案才會進行更新
    """
    stats = {"cached": 0, "generated": 0}
    progress_bar = EmojiProgressBar(len(source_files))  # 初始化進度條
    new_file_names = None
    if new_items is not None:
        new_file_names = [
            truncate_filename(
                append_export_extension_to_path(item["name"], item["mimeType"])
            )
            for item in new_items
        ]
    for index, (source_file, output_folder) in enumerate(
        zip(source_files, output_folders)
    ):
        if source_file.endswith("-metadata.json"):
            continue
        if new_file_names is not None:
            if os.path.basename(source_file) not in new_file_names:
                continue
        try:
            lst_rag_metadata = process_single_file(
                source_file,
                output_folder,
                chars_per_page,
                force,
                stats,
                gen_page_imgs,
            )
            # if google_drive_folder_id is not None and len(lst_rag_metadata) > 0:
            #     update_drive_folder_after_process_file(
            #         google_drive_folder_id, lst_rag_metadata
            #     )
            if metadata_dir and len(lst_rag_metadata) > 0:
                update_download_metadata_after_process_file(
                    metadata_dir, lst_rag_metadata
                )
            progress_bar.update(index + 1)  # 更新進度條
        except Exception as exc:
            import traceback

            traceback.print_exc()
            log_handler.write_log(f"{source_file} generated an exception: {exc}")
    print(
        f"split_txts.py, cached files: {stats['cached']}, generated files: {stats['generated']}"
    )


async def process_single_file_by_storage(
    drive_file: DriveFile,
    output_folder: str,
    chars_per_page: int,
    force: bool,
) -> List[RagMetadata]:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    # file_name, file_type = os.path.splitext(os.path.basename(source_file))
    file_name = drive_file.name
    file_type = drive_file.name.split(".")[-1]
    storage_client: StorageClient = storage_client_factory()

    # if (
    #     not force
    #     and await is_cached_by_storage(storage_client, output_folder)
    #     and stats is not None
    # ):
    #     stats["cached"] += 1
    #     return []

    lst_rag_metadata = []

    if file_type in [".docx", ".doc"]:
        lst_rag_metadata = await process_word_by_storage(
            chars_per_page,
            output_folder,
            drive_file,
            storage_client,
        )
    elif file_type == ".pdf":
        lst_rag_metadata = await process_pdf_by_storage(
            chars_per_page,
            file_type,
            output_folder,
            drive_file,
            storage_client,
        )
    elif file_type in [".ppt", ".pptx"]:
        lst_rag_metadata = await process_powerpoint_by_storage(
            chars_per_page,
            file_type,
            output_folder,
            drive_file,
            storage_client,
        )
    elif file_type in [".xls", ".xlsx"]:
        lst_rag_metadata = await process_excel_by_storage(
            chars_per_page,
            file_type,
            output_folder,
            drive_file,
            storage_client,
        )
    elif file_type == ".ods":
        converted_file = await convert_office_file_by_storage(
            drive_file, ".xlsx", storage_client
        )
        lst_rag_metadata = await process_excel_by_storage(
            converted_file, output_folder, file_name, ".xlsx", storage_client
        )
        await storage_client.delete_file(converted_file)
    elif file_type == ".odp":
        converted_file = await convert_office_file_by_storage(
            drive_file, ".pptx", storage_client
        )
        lst_rag_metadata = await process_powerpoint_by_storage(
            chars_per_page,
            file_name,
            file_type,
            output_folder,
            converted_file,
            storage_client,
        )
        await storage_client.delete_file(converted_file)
    elif file_type == ".rtf" or file_type == ".odt":
        converted_file = await convert_office_file_by_storage(
            drive_file, ".docx", storage_client
        )
        lst_rag_metadata = await process_word_by_storage(
            chars_per_page,
            output_folder,
            converted_file,
            storage_client,
        )
        await storage_client.delete_file(converted_file)
    elif file_type == ".txt" or await is_text_file_by_storage(
        drive_file, storage_client
    ):
        lst_rag_metadata = await process_txt_by_storage(
            chars_per_page,
            file_type,
            output_folder,
            drive_file,
            storage_client,
        )

    return lst_rag_metadata


async def is_cached_by_storage(
    storage_client: StorageClient, output_folder: str
) -> bool:
    # 實現檢查 storage 中是否已有緩存的邏輯
    return await storage_client.file_exists(f"{output_folder}/cache_flag")


async def convert_office_file_by_storage(
    drive_file: DriveFile, target_extension: str, storage_client: StorageClient
) -> DriveFile:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    # 創建臨時目錄
    with tempfile.TemporaryDirectory() as temp_dir:
        # 從 Google Storage 下載文件
        source_file_path = os.path.join(
            temp_dir, os.path.basename(drive_file.save_path)
        )
        file_content = await storage_client.retrieve_file(drive_file.save_path)

        # 將 BytesIO 內容寫入臨時文件
        with open(source_file_path, "wb") as f:
            f.write(file_content.getvalue())

        # 準備轉換後的文件名稱
        converted_file_base = os.path.splitext(os.path.basename(drive_file.save_path))[
            0
        ]
        converted_file_name = f"{converted_file_base}.{target_extension.lstrip('.')}"
        converted_file_path = os.path.join(temp_dir, converted_file_name)

        # 建構 LibreOffice 的命令列格式，用於檔案轉換
        libreoffice_command = os.getenv("LIBRE_OFFICE_COMMAND", "libreoffice")
        cmd = [
            libreoffice_command,
            "--headless",
            "--convert-to",
            target_extension.lstrip("."),
            "--outdir",
            temp_dir,
            source_file_path,
        ]

        try:
            # 設置 5 分鐘 (300 秒) 的超時限制
            proc = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                check=True,
                timeout=300,  # 5 分鐘超時
            )
            stdout, stderr = proc.stdout, proc.stderr
            log_handler.write_log(f"Stdout: {stdout}")
            log_handler.write_log(f"Stderr: {stderr}")
        except TimeoutExpired:
            error_message = (
                f"Conversion of {drive_file.save_path} timed out after 5 minutes"
            )
            log_handler.write_log(error_message)
            raise Exception(error_message)
        except subprocess.CalledProcessError as e:
            error_message = f"Error converting {drive_file.save_path}. Return code: {e.returncode}\nStdout: {e.stdout}\nStderr: {e.stderr}"
            log_handler.write_log(error_message)
            raise Exception(error_message)
        except Exception as e:
            log_handler.write_log(f"Error converting {drive_file.save_path}: {e}")
            raise Exception(f"Error converting {drive_file.save_path}: {e}")

        if not os.path.exists(converted_file_path):
            error_message = f"Converted file not found: {converted_file_path}"
            log_handler.write_log(error_message)
            raise Exception(error_message)

        # 上傳轉換後的文件到 Google Storage
        new_storage_path = os.path.join(
            os.path.dirname(drive_file.save_path), converted_file_name
        )
        with open(converted_file_path, "rb") as f:
            await storage_client.upload_file(new_storage_path, f)

        # 創建新的 DriveFile 對象，這個檔案不會更新到 firestore，所以重要的只有 save_path
        converted_drive_file = DriveFile(
            id=drive_file.id,
            name=converted_file_name,
            mimeType=drive_file.mimeType,
            modifiedTime=drive_file.modifiedTime,
            save_path=new_storage_path,
            folder_id=drive_file.folder_id,
            size=drive_file.size,
            parent=drive_file.parent,
            path=drive_file.path,
        )

        return converted_drive_file


async def is_text_file_by_storage(
    file_path: str, storage_client: StorageClient
) -> bool:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    try:
        file_content = await storage_client.get_file_content(file_path)
        file_content.decode("utf-8")
        return True
    except UnicodeDecodeError:
        return False


async def process_word_by_storage(
    chars_per_page,
    output_folder,
    drive_file: DriveFile,
    storage_client: StorageClient,
) -> List[RagMetadata]:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    pdf_drive_file = await convert_office_file_by_storage(
        drive_file, ".pdf", storage_client
    )
    lst_rag_metadata = await process_pdf_by_storage(
        chars_per_page,
        ".pdf",
        output_folder,
        pdf_drive_file,
        storage_client,
    )
    await storage_client.delete_file(pdf_drive_file.save_path)
    return lst_rag_metadata


async def process_pdf_by_storage(
    max_chars_per_file,
    file_type,
    output_folder,
    drive_file: DriveFile,
    storage_client: StorageClient,
) -> List[RagMetadata]:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    file_content = await storage_client.retrieve_file(drive_file.save_path)
    doc = fitz.open(stream=file_content, filetype="pdf")
    lst_rag_metadata = []

    for page_number in range(len(doc)):
        page = doc.load_page(page_number)
        text = page.get_text()
        paragraphs = text.split("\n")

        current_chunk = []
        current_char_count = 0
        chunk_number = 1

        for paragraph in paragraphs:
            if (
                current_char_count + len(paragraph) > max_chars_per_file
                and current_chunk
            ):
                await save_chunk_by_storage(
                    current_chunk,
                    output_folder,
                    drive_file.name,
                    f"{get_page_number_prefix()}{page_number + 1}",
                    f"{chunk_number:03d}",
                    storage_client,
                )
                lst_rag_metadata.append(
                    RagMetadata(
                        name=get_save_chunk_file_name(
                            drive_file.name,
                            f"{get_page_number_prefix()}{page_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                        ori_file_name=f"{drive_file.name}{file_type}",
                        page_number=page_number + 1,
                        save_path=get_save_chunk_file_path_by_storage(
                            output_folder,
                            drive_file.name,
                            f"{get_page_number_prefix()}{page_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                    )
                )
                current_chunk = []
                current_char_count = 0
                chunk_number += 1

            current_chunk.append(paragraph)
            current_char_count += len(paragraph)

        if current_chunk:
            await save_chunk_by_storage(
                current_chunk,
                output_folder,
                drive_file.name,
                f"{get_page_number_prefix()}{page_number + 1}",
                f"{chunk_number:03d}",
                storage_client,
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        drive_file.name,
                        f"{get_page_number_prefix()}{page_number + 1}",
                        f"{chunk_number:03d}",
                    ),
                    ori_file_name=f"{drive_file.name}{file_type}",
                    page_number=page_number + 1,
                    save_path=get_save_chunk_file_path_by_storage(
                        output_folder,
                        drive_file.name,
                        f"{get_page_number_prefix()}{page_number + 1}",
                        f"{chunk_number:03d}",
                    ),
                )
            )

    return lst_rag_metadata


async def process_powerpoint_by_storage(
    max_chars_per_file,
    file_type,
    output_folder,
    drive_file: DriveFile,
    storage_client: StorageClient,
) -> List[RagMetadata]:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    file_content = await storage_client.retrieve_file(drive_file.save_path)
    prs = Presentation(BytesIO(file_content))
    lst_rag_metadata = []

    for slide_number, slide in enumerate(prs.slides):
        text = extract_text_from_pptx(slide)
        paragraphs = text.split("\n")

        current_chunk = []
        current_char_count = 0
        chunk_number = 1

        for paragraph in paragraphs:
            if (
                current_char_count + len(paragraph) > max_chars_per_file
                and current_chunk
            ):
                await save_chunk_by_storage(
                    current_chunk,
                    output_folder,
                    drive_file.name,
                    f"{get_page_number_prefix()}{slide_number + 1}",
                    f"{chunk_number:03d}",
                    storage_client,
                )
                lst_rag_metadata.append(
                    RagMetadata(
                        name=get_save_chunk_file_name(
                            drive_file.name,
                            f"{get_page_number_prefix()}{slide_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                        ori_file_name=f"{drive_file.name}{file_type}",
                        page_number=slide_number + 1,
                        save_path=get_save_chunk_file_path_by_storage(
                            output_folder,
                            drive_file.name,
                            f"{get_page_number_prefix()}{slide_number + 1}",
                            f"{chunk_number:03d}",
                        ),
                    )
                )
                current_chunk = []
                current_char_count = 0
                chunk_number += 1

            current_chunk.append(paragraph)
            current_char_count += len(paragraph)

        if current_chunk:
            await save_chunk_by_storage(
                current_chunk,
                output_folder,
                drive_file.name,
                f"{get_page_number_prefix()}{slide_number + 1}",
                f"{chunk_number:03d}",
                storage_client,
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        drive_file.name,
                        f"{get_page_number_prefix()}{slide_number + 1}",
                        f"{chunk_number:03d}",
                    ),
                    ori_file_name=f"{drive_file.name}{file_type}",
                    page_number=slide_number + 1,
                    save_path=get_save_chunk_file_path_by_storage(
                        output_folder,
                        drive_file.name,
                        f"{get_page_number_prefix()}{slide_number + 1}",
                        f"{chunk_number:03d}",
                    ),
                )
            )

    return lst_rag_metadata


async def process_excel_by_storage(
    max_chars_per_file,
    file_name,
    file_type,
    output_folder,
    drive_file: DriveFile,
    storage_client: StorageClient,
) -> List[RagMetadata]:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    file_content = await storage_client.retrieve_file(drive_file.save_path)
    xls = pd.ExcelFile(BytesIO(file_content))
    lst_rag_metadata = []

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name)
        content_list = df.apply(lambda row: " ".join(row.astype(str)), axis=1).tolist()

        current_chunk = []
        current_char_count = 0
        file_counter = 1

        for row in content_list:
            if current_char_count + len(row) > max_chars_per_file and current_chunk:
                await save_chunk_by_storage(
                    current_chunk,
                    output_folder,
                    file_name,
                    f"{sheet_name}_{get_page_number_prefix()}_1_",
                    file_counter,
                    storage_client,
                )
                lst_rag_metadata.append(
                    RagMetadata(
                        name=get_save_chunk_file_name(
                            file_name,
                            f"{sheet_name}_{get_page_number_prefix()}_1_",
                            file_counter,
                        ),
                        ori_file_name=f"{file_name}{file_type}",
                        page_number=PAGE_NUMBER_NA,
                        sheet_name=sheet_name,
                        save_path=get_save_chunk_file_path_by_storage(
                            output_folder,
                            file_name,
                            f"{sheet_name}_{get_page_number_prefix()}_1_",
                            file_counter,
                        ),
                    )
                )
                current_chunk = []
                current_char_count = 0
                file_counter += 1

            current_chunk.append(row)
            current_char_count += len(row)

        if current_chunk:
            await save_chunk_by_storage(
                current_chunk,
                output_folder,
                file_name,
                f"{sheet_name}_{get_page_number_prefix()}_1_",
                file_counter,
                storage_client,
            )
            lst_rag_metadata.append(
                RagMetadata(
                    name=get_save_chunk_file_name(
                        file_name,
                        f"{sheet_name}_{get_page_number_prefix()}_1_",
                        file_counter,
                    ),
                    ori_file_name=f"{file_name}{file_type}",
                    page_number=PAGE_NUMBER_NA,
                    sheet_name=sheet_name,
                    save_path=get_save_chunk_file_path_by_storage(
                        output_folder,
                        file_name,
                        f"{sheet_name}_{get_page_number_prefix()}_1_",
                        file_counter,
                    ),
                )
            )

    return lst_rag_metadata


async def process_txt_by_storage(
    chars_per_page,
    file_type,
    output_folder,
    drive_file: DriveFile,
    storage_client: StorageClient,
) -> List[RagMetadata]:
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    # 實現處理文本文件並存儲到 storage 的邏輯
    pass


async def save_chunk_by_storage(
    chunk,
    output_folder,
    file_name,
    sheet_name,
    file_counter,
    storage_client: StorageClient,
):
    # 改到一半，還沒測試，但是換方向，所以 code 留著
    output_path = get_save_chunk_file_path_by_storage(
        output_folder, file_name, sheet_name, file_counter
    )
    content = "\n".join(chunk)
    await storage_client.store_file(output_path, BytesIO(content.encode("utf-8")))


if __name__ == "__main__":
    p = argparse.ArgumentParser(description="Split specified files into pages or parts")
    a = p.add_argument
    a("--source_files", nargs="+", help="Paths of the files to be split")
    a("--output_folders", nargs="+", help="Output folders for the split files")
    a("--chars_per_page", type=int, default=2000)
    a("--force", action="store_true")
    args = p.parse_args()

    split_txts_no_threads(
        args.source_files, args.output_folders, args.chars_per_page, args.force
    )
