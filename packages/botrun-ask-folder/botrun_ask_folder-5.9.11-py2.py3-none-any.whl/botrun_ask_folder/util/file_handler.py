from typing import List
import pandas as pd
import mimetypes
import os
import docx
from PyPDF2 import PdfReader
from pptx import Presentation
import pdfplumber

from botrun_ask_folder.split_txts import convert_office_file, extract_text_from_pptx


class UnsupportedFileException(Exception):
    """
    當檔案格式不支援時拋出的異常
    """

    def __init__(self, mime_type: str):
        self.mime_type = mime_type
        self.message = f"Unsupported file type: {mime_type}"
        super().__init__(self.message)


class HandlePowerpointError(Exception):
    def __init__(
        self,
    ):
        self.message = ""
        super().__init__(self.message)


class FailedToExtractContentException(Exception):
    """
    當無法從檔案中提取文字時拋出的異常
    """

    def __init__(self):
        self.message = "Failed to extract content from file."
        super().__init__(self.message)


def read_excel_file(file_path: str) -> str:
    """
    讀取 Excel 檔案的通用函數，會嘗試不同的讀取方式

    Args:
        file_path: Excel 檔案路徑

    Returns:
        pandas DataFrame 物件

    Raises:
        Exception: 當所有讀取方式都失敗時拋出
    """
    # 方法1: 直接讀取
    try:
        df = pd.read_excel(file_path)
        return df.to_string(index=False)
    except Exception as e:
        print(f"Standard read failed: {str(e)}")

    try:
        import pylightxl as xl

        db = xl.readxl(fn=file_path)

        # 取得第一個工作表
        ws = db.ws(db.ws_names[0])
        content = ""
        # 逐行讀取並輸出
        for row in ws.rows:
            # 將每個 cell 轉為字串並用逗號連接
            line = ",".join(str(cell) for cell in row)
            content += line + "\n"
        return content

    except Exception as e:
        print(f"Error: {str(e)}")
    # 方法4: 最保守的讀取方式
    return pd.read_excel(
        file_path,
        dtype=str,
        engine="openpyxl",
        keep_default_na=False,
    ).to_string(index=False)


async def handle_file_upload(
    file_name: str,
    file_path: str,
    file_mime: str,
) -> str:
    content = ""
    # 獲取實際的 MIME 類型
    mime_type, _ = mimetypes.guess_type(file_path)

    if mime_type == "application/pdf" or file_mime == "application/pdf":
        content = extract_text_from_pdf(file_path)
    elif mime_type in [
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ] or file_mime in [
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ]:
        # 如果是 .doc 文件，先轉換為 .docx
        if file_name.lower().endswith(".doc"):
            converted_file_path = convert_office_file(file_path, ".docx")
            content = extract_text_from_docx(converted_file_path)
            os.remove(converted_file_path)
        else:
            content = extract_text_from_docx(file_path)
    elif mime_type in [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ] or file_mime in [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ]:
        content = read_excel_file(file_path)
    elif mime_type in [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.oasis.opendocument.presentation",
    ] or file_mime in [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.oasis.opendocument.presentation",
    ]:
        try:
            # 總是先轉為 .pptx 格式
            converted_file_path = convert_office_file(file_path, ".pptx")
            prs = Presentation(converted_file_path)
            content = ""
            for slide in prs.slides:
                content += extract_text_from_pptx(slide) + "\n\n"
            os.remove(converted_file_path)
        except Exception as e:
            import traceback

            traceback.print_exc()

            print(f"Error processing PowerPoint file: {e}")
            raise HandlePowerpointError()
    elif (
        mime_type == "application/vnd.oasis.opendocument.spreadsheet"
        or file_mime == "application/vnd.oasis.opendocument.spreadsheet"
    ):
        converted_file_path = convert_office_file(file_path, ".xlsx")
        df = pd.read_excel(converted_file_path)
        content = df.to_string(index=False)
        os.remove(converted_file_path)
    elif (
        mime_type == "application/vnd.oasis.opendocument.presentation"
        or file_mime == "application/vnd.oasis.opendocument.presentation"
    ):
        converted_file_path = convert_office_file(file_path, ".pptx")
        prs = Presentation(converted_file_path)
        content = extract_text_from_pptx(prs)
        os.remove(converted_file_path)
    elif mime_type in [
        "application/rtf",
        "application/vnd.oasis.opendocument.text",
        "text/rtf",
    ] or file_mime in [
        "application/rtf",
        "application/vnd.oasis.opendocument.text",
        "text/rtf",
    ]:
        converted_file_path = convert_office_file(file_path, ".docx")
        content = extract_text_from_docx(converted_file_path)
        os.remove(converted_file_path)
    # elif mime_type == "text/plain" or file_mime == "text/plain":
    #     with open(file_path, "r", encoding="utf-8") as f:
    #         content = f.read()
    # else:
    #     raise UnsupportedFileException(mime_type or file_mime)
    else:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
                if not isinstance(content, str):
                    raise UnsupportedFileException(mime_type or file_mime)
        except Exception as e:
            print(f"Error reading file: {e}")
            raise UnsupportedFileException(mime_type or file_mime)

    if not content:
        raise FailedToExtractContentException()
    return content


def extract_text_from_pdf(file_path, table_settings=None):
    if table_settings is None:
        table_settings = {
            "vertical_strategy": "lines",
            "horizontal_strategy": "lines",
            "intersection_tolerance": 3,
        }

    text_content = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            # 獲取表格位置和內容
            tables = page.extract_tables(table_settings)
            table_bboxes = [table.bbox for table in page.find_tables()]

            # 使用 page.chars 而不是 extract_words 來處理 CID 編碼問題
            chars = page.chars

            # 按照位置排序字符
            chars.sort(key=lambda x: (-x["top"], x["x0"]))

            # 組織文字，處理 CID 編碼
            current_line = []
            current_y = None
            non_table_text = []

            for char in chars:
                # 檢查是否在表格內
                in_table = False
                char_bbox = (char["x0"], char["top"], char["x1"], char["bottom"])
                for table_bbox in table_bboxes:
                    if (
                        char_bbox[0] >= table_bbox[0]
                        and char_bbox[1] >= table_bbox[1]
                        and char_bbox[2] <= table_bbox[2]
                        and char_bbox[3] <= table_bbox[3]
                    ):
                        in_table = True
                        break

                if not in_table:
                    # 使用 get 方法安全地獲取文字，如果沒有則使用空字串
                    char_text = char.get("text", "")

                    # 如果是新的一行
                    if current_y is None or abs(char["top"] - current_y) > 3:
                        if current_line:
                            non_table_text.append("".join(current_line))
                        current_line = [char_text]
                        current_y = char["top"]
                    else:
                        current_line.append(char_text)

            # 加入最後一行
            if current_line:
                non_table_text.append("".join(current_line))

            # 加入非表格文字
            if non_table_text:
                text_content.append("\n".join(non_table_text))

            # 加入表格內容（使用 Markdown 格式）
            if tables:
                for table in tables:
                    text_content.append("\n")  # 表格前空一行
                    # 表頭
                    text_content.append(
                        "|"
                        + "|".join(f" {str(cell or '')} " for cell in table[0])
                        + "|"
                    )
                    # 分隔線
                    text_content.append("|" + "|".join("---" for _ in table[0]) + "|")
                    # 表格內容
                    for row in table[1:]:
                        text_content.append(
                            "|" + "|".join(f" {str(cell or '')} " for cell in row) + "|"
                        )
                    text_content.append("\n")  # 表格後空一行

            text_content.append("\n")  # 每頁後空一行

    return "\n".join(text_content)


def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    text = "\n".join(text)
    return text
