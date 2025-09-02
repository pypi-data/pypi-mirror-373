#!/usr/bin/env python3

import asyncio
import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
import pandas as pd
import pypdf
from pptx import Presentation
from docx import Document
import aiofiles
import io

logger = logging.getLogger(__name__)

class FileHandlerError(Exception):
    pass

class BaseFileHandler(ABC):
    
    @abstractmethod
    async def read_file(self, file_path: Path, **kwargs) -> str:
        pass
    
    def _validate_file(self, file_path: Path) -> None:
        if not file_path.exists():
            raise FileHandlerError(f"File does not exist: {file_path}")
        
        if not file_path.is_file():
            raise FileHandlerError(f"Not a file: {file_path}")
        
        if file_path.stat().st_size == 0:
            raise FileHandlerError(f"Empty file: {file_path}")

class ExcelHandler(BaseFileHandler):
    
    async def read_file(self, file_path: Path, sheet_name: Optional[str] = None, **kwargs) -> str:
        self._validate_file(file_path)
        
        try:
            loop = asyncio.get_event_loop()
            
            def read_excel():
                try:
                    if sheet_name:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                    else:
                        excel_file = pd.ExcelFile(file_path)
                        all_sheets = []
                        
                        for sheet in excel_file.sheet_names:
                            df = pd.read_excel(file_path, sheet_name=sheet)
                            all_sheets.append(f"=== Sheet: {sheet} ===\n{df.to_string()}\n")
                        
                        return "\n".join(all_sheets)
                    
                    return df.to_string()
                except Exception as e:
                    raise FileHandlerError(f"Excel file reading error: {str(e)}")
            
            content = await loop.run_in_executor(None, read_excel)
            return content
            
        except Exception as e:
            raise FileHandlerError(f"Excel file processing error: {str(e)}")

class PDFHandler(BaseFileHandler):
    
    async def read_file(self, file_path: Path, page_range: Optional[str] = None, **kwargs) -> str:
        self._validate_file(file_path)
        
        try:
            loop = asyncio.get_event_loop()
            
            def read_pdf():
                try:
                    with open(file_path, 'rb') as file:
                        pdf_reader = pypdf.PdfReader(file)
                        total_pages = len(pdf_reader.pages)
                        
                        pages_to_read = self._parse_page_range(page_range, total_pages)
                        
                        text_content = []
                        for page_num in pages_to_read:
                            if 0 <= page_num < total_pages:
                                page = pdf_reader.pages[page_num]
                                text = page.extract_text()
                                if text.strip():
                                    text_content.append(f"=== Page {page_num + 1} ===\n{text}\n")
                        
                        return "\n".join(text_content)
                except Exception as e:
                    raise FileHandlerError(f"PDF file reading error: {str(e)}")
            
            content = await loop.run_in_executor(None, read_pdf)
            return content
            
        except Exception as e:
            raise FileHandlerError(f"PDF file processing error: {str(e)}")
    
    def _parse_page_range(self, page_range: Optional[str], total_pages: int) -> List[int]:
        if not page_range:
            return list(range(total_pages))
        
        pages = []
        try:
            for part in page_range.split(','):
                part = part.strip()
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    pages.extend(range(start - 1, end))
                else:
                    pages.append(int(part) - 1)
        except ValueError:
            raise FileHandlerError(f"Invalid page range format: {page_range}")
        
        return pages

class PPTHandler(BaseFileHandler):
    
    async def read_file(self, file_path: Path, **kwargs) -> str:
        self._validate_file(file_path)
        
        try:
            loop = asyncio.get_event_loop()
            
            def read_ppt():
                try:
                    presentation = Presentation(file_path)
                    content = []
                    
                    for i, slide in enumerate(presentation.slides):
                        slide_content = [f"=== Slide {i + 1} ==="]
                        
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content.append(shape.text.strip())
                        
                        if len(slide_content) > 1:
                            content.append("\n".join(slide_content))
                    
                    return "\n\n".join(content)
                except Exception as e:
                    raise FileHandlerError(f"PowerPoint file reading error: {str(e)}")
            
            content = await loop.run_in_executor(None, read_ppt)
            return content
            
        except Exception as e:
            raise FileHandlerError(f"PowerPoint file processing error: {str(e)}")

class WordHandler(BaseFileHandler):
    
    async def read_file(self, file_path: Path, **kwargs) -> str:
        self._validate_file(file_path)
        
        try:
            loop = asyncio.get_event_loop()
            
            def read_word():
                try:
                    doc = Document(file_path)
                    content = []
                    
                    for paragraph in doc.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            content.append(text)
                    
                    for table in doc.tables:
                        table_content = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                table_content.append(" | ".join(row_text))
                        
                        if table_content:
                            content.append("=== Table ===\n" + "\n".join(table_content))
                    
                    return "\n".join(content)
                except Exception as e:
                    raise FileHandlerError(f"Word file reading error: {str(e)}")
            
            content = await loop.run_in_executor(None, read_word)
            return content
            
        except Exception as e:
            raise FileHandlerError(f"Word file processing error: {str(e)}")

class FileCache:
    
    def __init__(self, max_size: int = 100):
        self.cache: Dict[str, str] = {}
        self.max_size = max_size
        self.access_order: List[str] = []
    
    def get(self, file_path: str, file_mtime: float) -> Optional[str]:
        cache_key = f"{file_path}:{file_mtime}"
        if cache_key in self.cache:
            self.access_order.remove(cache_key)
            self.access_order.append(cache_key)
            return self.cache[cache_key]
        return None
    
    def set(self, file_path: str, file_mtime: float, content: str) -> None:
        cache_key = f"{file_path}:{file_mtime}"
        
        if len(self.cache) >= self.max_size:
            oldest_key = self.access_order.pop(0)
            del self.cache[oldest_key]
        
        self.cache[cache_key] = content
        self.access_order.append(cache_key)
    
    def clear(self) -> None:
        self.cache.clear()
        self.access_order.clear()

file_cache = FileCache()
